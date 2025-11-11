#!/usr/bin/env python3
"""
Generate TrueHome Pitch Deck PowerPoint Presentation
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_pitch_deck():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors (TrueHome brand colors - adjust as needed)
    primary_color = RGBColor(0, 102, 204)  # Blue
    accent_color = RGBColor(255, 140, 0)   # Orange
    text_color = RGBColor(51, 51, 51)      # Dark gray
    light_gray = RGBColor(128, 128, 128)
    
    # SLIDE 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "TrueHome"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Carfax for Homes - Know Your Home's Complete History"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = light_gray
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Pitch Deck 2024"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(14)
    footer_para.font.color.rgb = light_gray
    footer_para.alignment = PP_ALIGN.CENTER
    
    # SLIDE 2: The Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "The Problem"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Home Buyers Are Flying Blind"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "68% of home buyers discover major issues after purchase"
    p.font.size = Pt(20)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "$15,000 average cost of unexpected repairs in first year"
    p.font.size = Pt(20)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "• No centralized property history database"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "• Hidden issues (flooding, structural problems, unpermitted work)"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "• Incomplete information from sellers"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "• Expensive inspections don't reveal history"
    p.font.size = Pt(18)
    
    # SLIDE 3: The Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The Solution"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Complete Property History at Your Fingertips"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "✓ Sales history, work records, insurance claims"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "✓ Title issues, environmental assessments"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "✓ Maintenance records and warranties"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "✓ AI-powered risk assessment and recommendations"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "✓ TrueHome Certified professional network"
    p.font.size = Pt(18)
    
    # SLIDE 4: Market Opportunity
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Market Opportunity"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Massive Market, Growing Demand"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "TAM: $200B+ US real estate market"
    p.font.size = Pt(20)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "6M+ home sales annually"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "SAM: 2M+ home buyers per year seeking reports"
    p.font.size = Pt(20)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "$100M+ annual market opportunity"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "SOM: Year 3 target - 100,000 reports = $5M revenue"
    p.font.size = Pt(20)
    p.font.bold = True
    
    # SLIDE 5: Product Demo
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "How It Works"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Simple, Powerful, Comprehensive"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "1. Search property address"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "2. View instant property overview"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "3. Purchase full report ($49.99)"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "4. Get comprehensive analysis with recommendations"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Key Features: Risk scoring, Timeline, Certified professionals"
    p.font.size = Pt(18)
    p.font.italic = True
    
    # SLIDE 6: Business Model
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Business Model"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Multiple Revenue Streams, High Margins"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "1. Individual Reports: $49.99 per report"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "2. Subscriptions: $9.99-$99.99/month (planned)"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "3. Marketplace: 10-15% commission (planned)"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "4. API Licensing: B2B data access (planned)"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "5. White-Label: Enterprise solutions (planned)"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Unit Economics: 90%+ gross margin, $200+ LTV, <$50 CAC"
    p.font.size = Pt(18)
    p.font.bold = True
    
    # SLIDE 7: Traction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Traction"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Early Traction, Strong Foundation"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "✓ MVP launched and operational"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "✓ Full report generation system"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "✓ Sample data and comprehensive reports"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "✓ Professional website with sample report"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "✓ TrueHome Certified professional network concept"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Ready to scale with investment"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.italic = True
    
    # SLIDE 8: Competitive Landscape
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Competitive Landscape"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "First-Mover Advantage in a Fragmented Market"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Home Inspections: One-time, expensive, no history"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "Public Records: Scattered, incomplete, hard to access"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "Real Estate Sites: Basic info only, no comprehensive history"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "TrueHome: Complete history, AI insights, certified network"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = primary_color
    
    # SLIDE 9: Go-to-Market
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Go-to-Market Strategy"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Multi-Channel Growth Strategy"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "B2C Channels:"
    p.font.size = Pt(20)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• SEO-optimized property pages"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• Content marketing & social media"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• Referral program & agent partnerships"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "B2B Channels:"
    p.font.size = Pt(20)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• Real estate platform integrations (Zillow, Redfin)"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• Mortgage lender & insurance partnerships"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• White-label solutions"
    p.font.size = Pt(16)
    
    # SLIDE 10: Team
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Team"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Experienced Team Building the Future of Real Estate"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "[Add founder names, backgrounds, and key achievements]"
    p.font.size = Pt(18)
    p.font.italic = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Key Strengths:"
    p.font.size = Pt(20)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• Real estate industry expertise"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• Technology & product development"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• Business development & partnerships"
    p.font.size = Pt(16)
    
    # SLIDE 11: Financial Projections
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Financial Projections"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Path to $10M ARR in 3 Years"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Year 1: $500K revenue | 10,000 reports | 8,000 customers"
    p.font.size = Pt(18)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "Year 2: $2.5M revenue | 50,000 reports | 40,000 customers"
    p.font.size = Pt(18)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "Year 3: $10M revenue | 200,000 reports | 150,000 customers"
    p.font.size = Pt(18)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Key Assumptions:"
    p.font.size = Pt(18)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• 5% conversion rate | $50 CAC | 20% repeat rate | 70%+ gross margins"
    p.font.size = Pt(16)
    
    # SLIDE 12: The Ask
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The Ask"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Join Us in Transforming Real Estate Transparency"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Funding Ask: $2M Seed Round"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = primary_color
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Use of Funds:"
    p.font.size = Pt(20)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• 40% - Product development (AI, mobile app, API)"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• 30% - Marketing & customer acquisition"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• 20% - Data partnerships & infrastructure"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• 10% - Team expansion"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Milestones: 1,000 customers (M6), $1M ARR (M18), $5M ARR (M24)"
    p.font.size = Pt(16)
    p.font.italic = True
    
    # SLIDE 13: Vision
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Vision"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Building the Property Intelligence Platform"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Long-term Vision:"
    p.font.size = Pt(20)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• Every property has a TrueHome history"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• Industry standard for property transparency"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• Platform connecting homeowners, buyers, and professionals"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• AI-powered property insights and predictions"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Expansion: International markets, commercial properties, insurance/lending"
    p.font.size = Pt(16)
    p.font.italic = True
    
    # SLIDE 14: Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Let's Build the Future of Real Estate Together"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    title_para.alignment = PP_ALIGN.CENTER
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(2))
    content_frame = content_box.text_frame
    content_frame.text = "Key Takeaways:"
    p = content_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p = content_frame.add_paragraph()
    p.text = "• Massive market opportunity"
    p.font.size = Pt(16)
    p = content_frame.add_paragraph()
    p.text = "• First-mover advantage"
    p.font.size = Pt(16)
    p = content_frame.add_paragraph()
    p.text = "• Multiple revenue streams"
    p.font.size = Pt(16)
    p = content_frame.add_paragraph()
    p.text = "• Strong unit economics"
    p.font.size = Pt(16)
    
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = "[Your Email] | [Your Website] | [Your LinkedIn]"
    contact_para = contact_frame.paragraphs[0]
    contact_para.font.size = Pt(14)
    contact_para.font.color.rgb = light_gray
    contact_para.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    filename = "TrueHome_Pitch_Deck.pptx"
    prs.save(filename)
    print(f"✅ Pitch deck created: {filename}")
    print(f"📁 Location: {filename}")
    return filename

if __name__ == "__main__":
    try:
        create_pitch_deck()
    except ImportError:
        print("❌ Error: python-pptx not installed")
        print("📦 Install with: pip install python-pptx")
    except Exception as e:
        print(f"❌ Error creating PowerPoint: {e}")
        print("💡 Make sure python-pptx is installed: pip install python-pptx")

